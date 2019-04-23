import numpy as np
from scipy import fftpack
import os
from pptx import Presentation
from pptx.util import Inches,Pt
from PyQt4 import QtGui
import glob
import hyperspy.api as hspy

def dataOut(xaxis,data,filename):
    out = np.array([xaxis,data])
    out = out.T   
    with open(filename,'wb') as f:
        np.savetxt(f,out,delimiter = ' , ',fmt='%2e')  
    f.close()
    return

def processRDF(filename,binsize=1,scalefactor=10,binning=None):
    data = hspy.load(filename)
    data = data.inav[0]
    if binning:
        outshape = [data.data.shape[0]/binning,data.data.shape[1]/binning]
        xaxis,profile = RDF(data.rebin(outshape),binsize,scalefactor)
        return(xaxis,profile)
    xaxis,profile = RDF(data,binsize,scalefactor)
    return(xaxis,profile)
	
def azimuthalAverage(image, binsize=0.5):
    # Calculate the indices from the image
    y, x = np.indices(image.shape)
    center = np.array([(x.max()-x.min())/2.0, (y.max()-y.min())/2.0])
    #mask = np.ones(image.shape,dtype='bool')
    #weights = np.ones(image.shape)
    
    r = np.hypot(x - center[0], y - center[1])

    # the 'bins' as initially defined are lower/upper bounds for each bin
    # so that values will be in [lower,upper)  
    nbins = int(np.round(r.max() / binsize)+1)
    maxbin = nbins * binsize
    bins = np.linspace(0,maxbin,nbins+1)
    # but we're probably more interested in the bin centers than their left or right sides...
    bin_centers = (bins[1:]+bins[:-1])/2.0

    # how many per bin (i.e., histogram)?
    # there are never any in bin 0, because the lowest index returned by digitize is 1
    #nr = np.bincount(whichbin)[1:]
    #nr = np.histogram(r,bins)[0]

    # recall that bins are from 1 to nbins (which is expressed in array terms by arange(nbins)+1 or xrange(1,nbins+1) )
    # radial_prof.shape = bin_centers.shape
    radial_prof = np.histogram(r, bins, weights=(image))[0] / np.histogram(r, bins,weights=np.ones(np.shape(image)))[0]
    return bin_centers,radial_prof

def RDF(image,binsize,scalefactor=None):
    psd = np.abs(fftpack.fftshift(fftpack.fft2(image.data)))**2
    xaxis,profile = azimuthalAverage(psd,binsize)
    if scalefactor:
        scale = scalefactor*image.axes_manager[0].scale
        xaxis = xaxis/(scale*len(psd))
    return(xaxis,profile)
	
def insertImages(prs,index,filelist='None'):
    if filelist == None:    
        filelist = QtGui.QFileDialog.getOpenFileNames(None,'Select images for slide ','C:/TEMP')     
        os.chdir(os.path.dirname(filelist[0]))
    n=len(filelist)
    w=(13.33-(n+1)*0.25)/n
    if w < 2.5:
        w = 2.5
        margin = (13.333-(n*w/2)-(0.25*(n-1)))/2
        for i in range(0,len(filelist)):        
            if i < len(filelist)/2:
                prs.slides[index].shapes.add_picture(filelist[i],left=Inches(margin+(i*0.25)+i*w),width=Inches(w),top=Inches(1.75))
            else:
                j = i - int(len(filelist)/2)
                prs.slides[index].shapes.add_picture(filelist[i],left=Inches(margin+(j*0.25)+j*w),width=Inches(w),top=Inches(1.75+w+0.25))
        return()
    else:
        if w > 5.5:
            w=5.5
    margin = (13.333-(n*w)-(0.25*(n-1)))/2
    for i in range(0,len(filelist)):
        prs.slides[index].shapes.add_picture(filelist[i],left=Inches(margin+(i*0.25)+i*w),width=Inches(w),top=Inches(1.75))
    return()
	
def makeSlide(prs,title,subtitle,index,filelist):
    ## Widescreen slide layout (13.333" X 7.5 in")
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    txBox = slide.shapes.add_textbox(Inches(0),Inches(0),Inches(10),Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]    
    p.text = title
    p.font.name = 'Calibri Light'
    p.font.size = Pt(60)
    
    txBox2 = slide.shapes.add_textbox(Inches(0),Inches(1),Inches(10),Inches(0.65))
    tf2 = txBox2.text_frame    
    p2 = tf2.paragraphs[0]    
    p2.text = subtitle
    p2.font.name = 'Calibri'
    p2.font.size = Pt(24)
    insertImages(prs,index,filelist)
    return()

def makePres(width=13.333,height=7.5):
    prs=Presentation()
    prs.slide_width = 72*Pt(width)#12192031
    prs.slide_height = 72*Pt(height)#6858189
    return(prs)
